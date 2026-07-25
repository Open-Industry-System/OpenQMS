import uuid
from datetime import UTC, date, datetime

import pytest
from pptx import Presentation  # python-pptx 读回断言

from app.models.capa import CAPAEightD, CapaRootCauseVerification
from app.models.supplier import Supplier
from app.models.supplier_risk import SupplierRiskAlert
from app.services import capa_ppt_service

pytestmark = pytest.mark.requires_db


async def _make_capa(db, factory_id, user_id, doc_no="8D-PPT-001"):
    capa = CAPAEightD(
        report_id=uuid.uuid4(), document_no=doc_no, title="8D 测试报告",
        product_line_code="DC-DC-100", factory_id=factory_id, created_by=user_id,
        status="D8_CLOSURE",
        d2_description="问题描述文本", d3_interim="遏制措施", d4_root_cause="根因",
        d5_correction="永久措施", d6_verification="验证", d7_prevention="预防", d8_closure="关闭",
        d1_team=[{"name": "张三", "role": "负责人"}],
    )
    db.add(capa)
    await db.flush()
    return capa


async def test_generate_content_returns_11_pages(db, admin_user, default_factory):
    capa = await _make_capa(db, default_factory.id, admin_user.user_id)
    content = await capa_ppt_service.generate_content(db, capa.report_id)
    assert len(content.pages) == 11
    assert content.pages[0].title == "封面"
    assert content.pages[2].title == "D2 问题描述"
    # D4 含 root_cause_verifications 字段（空列表，因无验证记录）
    assert content.root_cause_verifications == []


async def test_generate_content_with_verification_and_evidence(db, admin_user, default_factory):
    capa = await _make_capa(db, default_factory.id, admin_user.user_id)
    db.add(CapaRootCauseVerification(
        verification_id=uuid.uuid4(), capa_id=capa.report_id, factory_id=default_factory.id,
        root_cause_text="根因A", method="reproduction", result="复现成功", is_verified=True,
        conclusion="passed",
        evidence_attachments=[{"filename": "evidence.png", "content_type": "image/png"}],
    ))
    await db.flush()
    content = await capa_ppt_service.generate_content(db, capa.report_id)
    assert len(content.root_cause_verifications) == 1
    ev = content.root_cause_verifications[0]
    assert ev["evidence_attachments"] == [{"filename": "evidence.png", "content_type": "image/png"}]


async def test_generate_content_with_linked_risk_alert(db, admin_user, default_factory):
    """联动附录含供应商风险预警（模型 SupplierRiskAlert 在 app.models.supplier_risk）。"""
    capa = await _make_capa(db, default_factory.id, admin_user.user_id)
    # 需先建一个 supplier（FK 约束；NOT NULL: supplier_no/factory_id/name/short_name/created_by）
    sup = Supplier(supplier_id=uuid.uuid4(), supplier_no="S001", name="S",
                   short_name="S", factory_id=default_factory.id, created_by=admin_user.user_id)
    db.add(sup)
    await db.flush()
    alert_id = uuid.uuid4()
    db.add(SupplierRiskAlert(
        alert_id=alert_id, supplier_id=sup.supplier_id, factory_id=default_factory.id,
        risk_level="high", risk_score=80, quality_score=70, delivery_score=60, compliance_score=90,
        rule_results={}, alert_type="initial", status="open",
        linked_capa_id=capa.report_id, snapshot_date=date(2026, 7, 8),
    ))
    await db.flush()
    content = await capa_ppt_service.generate_content(db, capa.report_id)
    assert len(content.linked_risk_alerts) == 1
    alert = content.linked_risk_alerts[0]
    assert alert["risk_level"] == "high"
    # 故事 §40/§71 要求「预警单号与状态」——必须含标识号
    assert alert["alert_id"] == str(alert_id)


async def test_d1_team_renders_member_names_not_dict_repr(db, admin_user, default_factory):
    """D1 团队页应渲染成员姓名/角色可读文本，而非 str(dict) 的 {'name':...} 字典字面量。"""
    capa = await _make_capa(db, default_factory.id, admin_user.user_id)
    capa.d1_team = [{"name": "张三", "role": "负责人"}, {"name": "李四", "role": "工程师"}]
    await db.flush()
    content = await capa_ppt_service.generate_content(db, capa.report_id)
    d1_values = [s["value"] for s in content.pages[1].sections]
    joined = " ".join(d1_values)
    assert "张三" in joined and "李四" in joined
    assert "{'name'" not in joined and "\"name\"" not in joined  # 不应出现 dict 字面量


async def test_render_pptx_returns_valid_pptx(db, admin_user, default_factory):
    capa = await _make_capa(db, default_factory.id, admin_user.user_id)
    content = await capa_ppt_service.generate_content(db, capa.report_id)
    meta = capa_ppt_service.ExportMeta(
        export_id=uuid.uuid4(), version="20260708T143022Z",
        generated_at=datetime.now(UTC), generated_by=admin_user.user_id,
    )
    pptx_bytes = capa_ppt_service.render_pptx(content, meta, "passed", 1)
    assert isinstance(pptx_bytes, bytes) and len(pptx_bytes) > 0
    # 读回断言 11 页
    import io
    prs = Presentation(io.BytesIO(pptx_bytes))
    assert len(prs.slides) == 11
    # 每页使用 Blank 布局（slide_layouts[6]）——无标题占位符（避免与手绘标题文本框重复）
    for slide in prs.slides:
        assert list(slide.placeholders) == [], "Blank 布局不应带占位符（Title Only 会带空标题占位符）"


async def test_validate_ppt_content_missing_page(db, admin_user, default_factory):
    capa = await _make_capa(db, default_factory.id, admin_user.user_id)
    content = await capa_ppt_service.generate_content(db, capa.report_id)
    content.pages = content.pages[:10]  # 破坏：缺一页
    issues = capa_ppt_service._validate_ppt_content(content, capa)
    assert any("11" in i or "页数" in i for i in issues)


async def test_validate_ppt_content_complete(db, admin_user, default_factory):
    capa = await _make_capa(db, default_factory.id, admin_user.user_id)
    content = await capa_ppt_service.generate_content(db, capa.report_id)
    issues = capa_ppt_service._validate_ppt_content(content, capa)
    assert issues == []


async def test_validate_ppt_content_empty_d3_value(db, admin_user, default_factory):
    """d3_interim 空字符串 → D3 页 section value 空白 → 应返回 issue（非仅查 sections 列表）。"""
    capa = await _make_capa(db, default_factory.id, admin_user.user_id)
    capa.d3_interim = ""  # _make_capa 默认 d3_interim="遏制措施"，置空
    await db.flush()
    content = await capa_ppt_service.generate_content(db, capa.report_id)
    # D3 页 sections 非空（含 {"label":"措施","value":""}），但 value 空白
    assert content.pages[3].sections, "sanity: sections list is non-empty"
    issues = capa_ppt_service._validate_ppt_content(content, capa)
    assert any("D3" in i for i in issues)


async def test_validate_ppt_content_empty_d1_team(db, admin_user, default_factory):
    """d1_team 空列表 → D1 页 sections=[] → 应返回 issue（8D 语义：关闭报告须有团队，与 seed 审查标准 D1-D8 非空一致）。"""
    capa = await _make_capa(db, default_factory.id, admin_user.user_id)
    capa.d1_team = []  # _make_capa 默认 d1_team=[{"name":"张三","role":"负责人"}]，置空
    await db.flush()
    content = await capa_ppt_service.generate_content(db, capa.report_id)
    # D1 页（idx 1）sections 为空列表
    assert content.pages[1].sections == [], "sanity: D1 sections empty"
    issues = capa_ppt_service._validate_ppt_content(content, capa)
    assert any("D1" in i for i in issues)
