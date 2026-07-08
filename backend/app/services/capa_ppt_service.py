"""8D 报告 PPT 生成服务（US-E2E-01.10）。

generate_content 组装 PptContent（不渲染 pptx），render_pptx 一次性渲染。
审查只操作 PptContent；最终 pptx 在 API 层审查后只渲染一次。
"""
from __future__ import annotations

import uuid
from dataclasses import dataclass
from datetime import datetime
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from app.models.capa import (
    CapaD7NodeAction,
    CAPAEightD,
    CapaPptExport,
    CapaRootCauseVerification,
)


@dataclass
class PptPage:
    title: str
    sections: list[dict]  # [{label, value}]


@dataclass
class PptContent:
    capa_id: uuid.UUID
    pages: list[PptPage]
    linked_fmea_node: dict | None
    linked_scars: list[dict]
    linked_risk_alerts: list[dict]
    root_cause_verifications: list[dict]  # 每条含 evidence_attachments


@dataclass
class ExportMeta:
    export_id: uuid.UUID
    version: str
    generated_at: datetime
    generated_by: uuid.UUID


_PAGE_TITLES = [
    "封面", "D1 团队", "D2 问题描述", "D3 遏制措施", "D4 根因分析",
    "D5 永久措施", "D6 实施验证", "D7 预防复发", "D8 关闭结论", "联动附录", "生成信息",
]


async def generate_content(db: AsyncSession, capa_id: uuid.UUID) -> PptContent:
    """从落库数据组装 PptContent（不渲染 pptx）。"""
    capa = await db.get(CAPAEightD, capa_id)
    if capa is None:
        raise ValueError("CAPA 不存在")

    # D4 验证记录（含 evidence_attachments）
    verifs = (await db.execute(select(CapaRootCauseVerification).where(
        CapaRootCauseVerification.capa_id == capa_id
    ))).scalars().all()
    root_cause_verifications = [
        {
            "root_cause_text": v.root_cause_text,
            "method": v.method,
            "result": v.result,
            "is_verified": v.is_verified,
            "evidence_attachments": v.evidence_attachments or [],
        }
        for v in verifs
    ]

    # D7 node-actions
    d7_actions = (await db.execute(select(CapaD7NodeAction).where(
        CapaD7NodeAction.capa_id == capa_id
    ))).scalars().all()
    d7_list = [
        {"failure_mode": a.failure_mode_node_id, "action": a.action,
         "prevention_control": a.prevention_control_name_after}
        for a in d7_actions
    ]

    # 关联 FMEA 节点详情
    linked_fmea_node = await _load_linked_fmea_node(db, capa)

    # 关联 SCAR / 供应商风险预警
    linked_scars = await _load_linked_scars(db, capa)
    linked_risk_alerts = await _load_linked_risk_alerts(db, capa)

    pages = [
        PptPage("封面", [
            {"label": "8D 单号", "value": capa.document_no},
            {"label": "标题", "value": capa.title},
            {"label": "严重度", "value": capa.severity},
            {"label": "产品线", "value": capa.product_line_code},
            {"label": "发起人", "value": str(capa.created_by)},
            {"label": "状态", "value": capa.status},
            {"label": "日期", "value": str(capa.created_at)},
        ]),
        PptPage("D1 团队", [{"label": "成员", "value": str(m)} for m in (capa.d1_team or [])]),
        PptPage("D2 问题描述", [{"label": "描述", "value": capa.d2_description or ""}]),
        PptPage("D3 遏制措施", [{"label": "措施", "value": capa.d3_interim or ""}]),
        PptPage("D4 根因分析", [
            {"label": "根因", "value": capa.d4_root_cause or ""},
            *[{"label": f"验证{i+1}", "value": f"{v['method']}/{v['result']}/verified={v['is_verified']}/附件={[a['filename'] for a in v['evidence_attachments']]}"}
              for i, v in enumerate(root_cause_verifications)],
        ]),
        PptPage("D5 永久措施", [{"label": "措施", "value": capa.d5_correction or ""}]),
        PptPage("D6 实施验证", [{"label": "验证", "value": capa.d6_verification or ""}]),
        PptPage("D7 预防复发", [
            {"label": "预防", "value": capa.d7_prevention or ""},
            *[{"label": f"node-action{i+1}", "value": str(a)} for i, a in enumerate(d7_list)],
        ]),
        PptPage("D8 关闭结论", [{"label": "结论", "value": capa.d8_closure or ""}]),
        PptPage("联动附录", [
            {"label": "关联 FMEA 节点", "value": str(linked_fmea_node) if linked_fmea_node else "无"},
            *[{"label": f"SCAR{i+1}", "value": str(s)} for i, s in enumerate(linked_scars)],
            *[{"label": f"风险预警{i+1}", "value": str(r)} for i, r in enumerate(linked_risk_alerts)],
        ]),
        PptPage("生成信息", []),  # 占位，render_pptx 时用 meta + review 填充
    ]
    return PptContent(
        capa_id=capa_id, pages=pages,
        linked_fmea_node=linked_fmea_node, linked_scars=linked_scars,
        linked_risk_alerts=linked_risk_alerts, root_cause_verifications=root_cause_verifications,
    )


async def _load_linked_fmea_node(db: AsyncSession, capa) -> dict | None:
    """从 fmea_ref_id + fmea_node_id 提取关联失效模式节点详情。"""
    if capa.fmea_ref_id is None or capa.fmea_node_id is None:
        return None
    from app.models.fmea import FMEADocument
    fmea = await db.get(FMEADocument, capa.fmea_ref_id)
    if fmea is None:
        return {"deleted": True, "fmea_ref_id": str(capa.fmea_ref_id)}
    # 从 graph_data 提取节点（简单实现：遍历 nodes 找匹配 id）
    graph = fmea.graph_data or {}
    for node in graph.get("nodes", []):
        if node.get("id") == capa.fmea_node_id:
            return {"fmea_document_no": fmea.document_no, "node": node}
    return {"fmea_document_no": fmea.document_no, "node_id": capa.fmea_node_id, "found": False}


async def _load_linked_scars(db: AsyncSession, capa) -> list[dict]:
    """反查 SupplierSCAR.capa_ref_id == capa_id（模型在 app.models.supplier:178）。"""
    from app.models.supplier import SupplierSCAR
    rows = (await db.execute(select(SupplierSCAR).where(SupplierSCAR.capa_ref_id == capa.report_id))).scalars().all()
    return [{"scar_no": r.scar_no, "status": r.status, "description": r.description} for r in rows]


async def _load_linked_risk_alerts(db: AsyncSession, capa) -> list[dict]:
    """反查 SupplierRiskAlert.linked_capa_id == capa_id（模型在 app.models.supplier_risk）。"""
    from app.models.supplier_risk import SupplierRiskAlert
    rows = (await db.execute(select(SupplierRiskAlert).where(SupplierRiskAlert.linked_capa_id == capa.report_id))).scalars().all()
    return [{"risk_level": r.risk_level, "status": r.status} for r in rows]


def render_pptx(
    content: PptContent, meta: ExportMeta, review_status: str, review_rounds: int,
) -> bytes:
    """从 PptContent + meta + review 一次性渲染最终 pptx。第 11 页用已知 review_status 直接渲染。"""
    prs = Presentation()
    # 使用默认幻灯片布局（空白）
    for page in content.pages:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # 5 = blank
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        title_box.text_frame.text = page.title
        title_box.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5))
        tf = body_box.text_frame
        for i, sec in enumerate(page.sections):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"{sec['label']}: {sec['value']}"
            p.font.size = Pt(14)
        # 第 11 页「生成信息」用 meta + review 填充
        if page.title == "生成信息":
            p = tf.add_paragraph()
            p.text = f"版本: {meta.version} | 审查状态: {review_status} | 审查轮数: {review_rounds}"
            p.font.size = Pt(12)
    output = BytesIO()
    prs.save(output)
    return output.getvalue()


def _validate_ppt_content(content: PptContent, capa) -> list[str]:
    """内置规则校验：页数=11、各 D 步必填非空、联动附录一致。不调 LLM。"""
    issues = []
    if len(content.pages) != 11:
        issues.append(f"页数应为 11，实际 {len(content.pages)}")
    # D1-D8 各页 section value 非空（与 seed 审查标准 "D1-D8 各页非空" 一致）
    # D1：capa.d1_team or [] 空列表 → sections=[] → issue（8D 语义：关闭报告须有团队）
    # D2-D8：generate_content 用 `capa.d3_interim or ""` 生成 {"label":"措施","value":""}，
    # sections 列表非空但 value 空白 → 必须查 value.strip()，否则空内容会通过校验。
    for idx, label in [(1, "D1"), (2, "D2"), (3, "D3"), (4, "D4"), (5, "D5"), (6, "D6"), (7, "D7"), (8, "D8")]:
        if idx < len(content.pages):
            page = content.pages[idx]
            if not page.sections or all(
                not str(s.get("value", "")).strip() for s in page.sections
            ):
                issues.append(f"{label} 页内容为空")
    # 联动附录：有 fmea_ref_id 但 linked_fmea_node 为 None
    if capa.fmea_ref_id is not None and content.linked_fmea_node is None:
        issues.append("有 fmea_ref_id 但联动附录未含关联 FMEA 节点")
    return issues


async def get_export(db: AsyncSession, export_id: uuid.UUID, capa_id: uuid.UUID) -> CapaPptExport | None:
    """查 PPT 生成记录（供 GET /ppt-exports/{id} 回读审查报告）。"""
    return (await db.execute(select(CapaPptExport).where(
        CapaPptExport.export_id == export_id, CapaPptExport.capa_id == capa_id,
    ))).scalar_one_or_none()
